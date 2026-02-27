"""파이프라인 통합 테스트 스크립트.

test_pptx_replace.py가 먼저 실행되어야 합니다 (더미 PPTX 생성).

테스트 모드:
  python test_pipeline.py --test fixed     # 고정 페이지만 (Gamma 불필요)
  python test_pipeline.py --test full --gamma-id <ID>  # 전체 (Gamma 포함)

사용법:
  1. test_pptx_replace.py 실행 (더미 템플릿 생성)
  2. python test_pipeline.py --test fixed
"""

import argparse
import sys
from pathlib import Path

from pptx import Presentation

from proposal_pipeline.config import GAMMA_API_KEY, GAMMA_BASE_URL, POLL_INTERVAL_SEC, POLL_TIMEOUT_SEC
from proposal_pipeline.gamma.client import GammaHttpClient
from proposal_pipeline.pptx_replacer.replacer import PptxFileReplacer
from proposal_pipeline.pptx_merger.merger import PptxFileMerger
from proposal_pipeline.pipeline.pipeline import DefaultProposalPipeline
from proposal_pipeline.pipeline.models import PipelineConfig, FixedPage, DynamicPage


TEMPLATE_DIR = Path(__file__).parent.parent / "templates" / "dummy"
OUTPUT_DIR = Path(__file__).parent.parent / "output_test" / "pipeline"

SAMPLE_VARIABLES = {
    "고객명": "현대자동차",
    "날짜": "2026-02-25",
    "담당자명": "박영희",
    "연락처": "02-1234-5678",
    "제안일": "2026-02-25",
    "금액": "50,000,000",
    "페인포인트": "반복 업무에 시간 소모가 크고, 데이터 분석 역량이 부족합니다.",
}


def test_fixed_only():
    """고정 페이지만으로 파이프라인 테스트 (Gamma 미사용)."""
    print("=== 파이프라인 테스트: 고정 페이지만 ===")
    print()

    # 템플릿 존재 확인
    templates = ["표지.pptx", "회사소개.pptx", "약관.pptx"]
    for t in templates:
        if not (TEMPLATE_DIR / t).exists():
            print(f"에러: {TEMPLATE_DIR / t} 없음. test_pptx_replace.py를 먼저 실행하세요.")
            return False

    # 설정
    config = PipelineConfig(
        fixedPages=[
            FixedPage(name="표지", templatePath=TEMPLATE_DIR / "표지.pptx", order=10),
            FixedPage(name="회사소개", templatePath=TEMPLATE_DIR / "회사소개.pptx", order=20),
            FixedPage(name="약관", templatePath=TEMPLATE_DIR / "약관.pptx", order=30),
        ],
        dynamicPages=[],
        outputDir=OUTPUT_DIR,
    )

    # 파이프라인 실행 (Gamma client는 더미 — 동적 페이지 없으므로 호출 안 됨)
    gamma_client = GammaHttpClient(api_key="dummy", base_url="https://dummy.test")
    replacer = PptxFileReplacer()
    merger = PptxFileMerger()

    pipeline = DefaultProposalPipeline(
        gamma_client=gamma_client,
        replacer=replacer,
        merger=merger,
    )

    print(f"변수: {SAMPLE_VARIABLES}")
    print(f"고정 페이지: {[fp.name for fp in config.fixed_pages]}")
    print()

    result = pipeline.run(SAMPLE_VARIABLES, config)

    print(f"출력 파일: {result.output_path}")
    print(f"고정 페이지: {result.fixed_count}")
    print(f"동적 페이지: {result.dynamic_count}")
    print(f"총 슬라이드: {result.total_slides}")
    print(f"에러: {result.errors}")
    print()

    # 검증
    all_ok = True

    # 1) 파일 존재
    if not result.output_path.exists():
        print("!! 실패: 출력 파일 없음")
        return False
    print("OK - 출력 파일 생성됨")

    # 2) 에러 없음
    if not result.is_success:
        print(f"!! 에러 발생: {result.errors}")
        all_ok = False
    else:
        print("OK - 에러 없음")

    # 3) 슬라이드 수
    if result.total_slides != 3:
        print(f"!! 슬라이드 수: {result.total_slides} (기대: 3)")
        all_ok = False
    else:
        print("OK - 슬라이드 수 정상 (3)")

    # 4) 치환 확인: 첫 슬라이드에서 고객명 확인
    prs = Presentation(str(result.output_path))
    slide1_text = ""
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                slide1_text += "".join(run.text for run in para.runs)

    if "현대자동차" in slide1_text:
        print("OK - 치환 확인 (현대자동차)")
    else:
        print(f"!! 치환 확인 실패 — 슬라이드 1 텍스트: {slide1_text[:80]}")
        all_ok = False

    # 5) 파일명 확인
    expected_name = "[팀스파르타] 현대자동차 제안서.pptx"
    if result.output_path.name == expected_name:
        print(f"OK - 파일명: {expected_name}")
    else:
        print(f"!! 파일명: {result.output_path.name} (기대: {expected_name})")
        all_ok = False

    # 6) 순서 확인: 마지막 슬라이드에 약관 관련 내용
    last_text = ""
    for shape in prs.slides[-1].shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                last_text += "".join(run.text for run in para.runs)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    last_text += cell.text

    if "제안" in last_text or "견적" in last_text or "유효" in last_text:
        print("OK - 순서 확인 (마지막=약관)")
    else:
        print(f"!! 순서 확인 실패")
        all_ok = False

    print()
    return all_ok


def test_full(gamma_id: str = None):
    """전체 파이프라인 테스트 (Gamma API 포함).

    gamma_id가 있으면 Template API, 없으면 Generate API 사용.
    """
    print("=== 파이프라인 테스트: 전체 (Gamma 포함) ===")
    print()

    if gamma_id:
        dynamic_page = DynamicPage(
            name="커리큘럼",
            mode="template",
            gammaId=gamma_id,
            promptTemplate=(
                "고객사: {{고객명}}\n"
                "페인포인트: {{페인포인트}}\n"
                "AI 업무 자동화 교육 커리큘럼을 제안해주세요.\n"
                "대상 기업의 페인포인트를 해결하는 방향으로 구성해주세요."
            ),
            order=20,
            exportAs="pptx",
        )
    else:
        dynamic_page = DynamicPage(
            name="커리큘럼",
            mode="generate",
            promptTemplate=(
                "{{고객명}}을 위한 AI 업무 자동화 교육 제안\n\n"
                "페인포인트: {{페인포인트}}\n\n"
                "위 페인포인트를 해결하기 위한 AI 업무 자동화 교육 커리큘럼을 구성해주세요.\n"
                "교육 목표, 주요 커리큘럼(세부 모듈), 기대 효과를 포함해주세요."
            ),
            order=20,
            numCards=5,
            exportAs="pptx",
        )

    config = PipelineConfig(
        fixedPages=[
            FixedPage(name="표지", templatePath=TEMPLATE_DIR / "표지.pptx", order=10),
            FixedPage(name="회사소개", templatePath=TEMPLATE_DIR / "회사소개.pptx", order=30),
            FixedPage(name="약관", templatePath=TEMPLATE_DIR / "약관.pptx", order=40),
        ],
        dynamicPages=[dynamic_page],
        outputDir=OUTPUT_DIR,
        gammaThemeId="jdvmtofxo715647",
    )

    gamma_client = GammaHttpClient(api_key=GAMMA_API_KEY, base_url=GAMMA_BASE_URL)
    replacer = PptxFileReplacer()
    merger = PptxFileMerger()

    pipeline = DefaultProposalPipeline(
        gamma_client=gamma_client,
        replacer=replacer,
        merger=merger,
    )

    print(f"변수: {SAMPLE_VARIABLES}")
    print(f"고정 페이지: {[fp.name for fp in config.fixed_pages]}")
    mode_str = f"template (gammaId={gamma_id})" if gamma_id else "generate"
    print(f"동적 페이지: {[dp.name for dp in config.dynamic_pages]} ({mode_str})")
    print()

    print("파이프라인 실행 중 (Gamma 생성 대기 포함)...")
    result = pipeline.run(SAMPLE_VARIABLES, config)

    print()
    print(f"출력 파일: {result.output_path}")
    print(f"고정 페이지: {result.fixed_count}")
    print(f"동적 페이지: {result.dynamic_count}")
    print(f"총 슬라이드: {result.total_slides}")
    print(f"에러: {result.errors}")
    print()

    if result.is_success:
        print("전체 파이프라인 테스트 성공!")
    else:
        print(f"에러 발생: {result.errors}")

    return result.is_success


def main():
    parser = argparse.ArgumentParser(description="파이프라인 통합 테스트")
    parser.add_argument("--test", required=True, choices=["fixed", "full"])
    parser.add_argument("--gamma-id", help="full 테스트 시 Gamma 템플릿 ID")
    args = parser.parse_args()

    if args.test == "fixed":
        ok = test_fixed_only()
    elif args.test == "full":
        if not GAMMA_API_KEY:
            print("에러: GAMMA_API_KEY가 설정되지 않았습니다.")
            sys.exit(1)
        ok = test_full(gamma_id=args.gamma_id)

    print()
    if ok:
        print("=== 테스트 통과! ===")
    else:
        print("=== 테스트 실패 ===")
        sys.exit(1)


if __name__ == "__main__":
    main()
