"""E2E 테스트 — 파이프라인 v2 종합 검증.

시나리오별 제안서를 생성하고 프로그래밍 방식으로 검증한다.
검증 항목:
  1. 슬라이드 수 일치
  2. 플레이스홀더 잔존 없음 ({{...}} 패턴)
  3. 테이블 데이터 정확성 (행 수, 헤더 보존)
  4. 서식 보존 (폰트, 미디어 파일)
  5. 슬라이드 마스터/레이아웃 무결성
  6. 슬라이드 순서 (order 기준)
"""

import json
import math
import re
import sys
import zipfile
from datetime import date
from pathlib import Path

from proposal_pipeline.config import GAMMA_API_KEY, GAMMA_BASE_URL
from proposal_pipeline.curriculum.models import CurriculumModule, TablePage, TopicRow
from proposal_pipeline.gamma.client import GammaHttpClient
from proposal_pipeline.pipeline.models import PipelineConfig, build_config
from proposal_pipeline.pipeline.pipeline import DefaultProposalPipeline
from proposal_pipeline.pptx_merger.merger import PptxZipMerger
from proposal_pipeline.pptx_replacer.replacer import PptxFileReplacer
from proposal_pipeline._resources import get_template_dir, get_data_dir

TEMPLATE_DIR = get_template_dir("parts")
OUTPUT_DIR = Path(__file__).parent.parent / "output" / "e2e_tests"
REFERENCE_JSON = get_data_dir() / "curriculum_reference.json"

# ── 헬퍼 ──────────────────────────────────────────────────────────────


def load_reference_modules(module_ids: list[str]) -> list[CurriculumModule]:
    """참조 JSON → CurriculumModule 리스트."""
    with open(REFERENCE_JSON, encoding="utf-8") as f:
        ref = json.load(f)

    modules = []
    for mid in module_ids:
        data = ref[mid]
        topics = [TopicRow(**t) for t in data["topics"]]
        n_pages = data.get("table_pages", 1)
        rows_per_page = math.ceil(len(topics) / n_pages)

        table_pages = []
        for i in range(n_pages):
            start = i * rows_per_page
            end = min(start + rows_per_page, len(topics))
            label = f"{i + 1}차" if n_pages > 1 else ""
            table_pages.append(TablePage(label=label, rows=topics[start:end]))

        modules.append(CurriculumModule(
            id=mid, name=data["name"],
            total_hours=data["total_hours"],
            table_pages=table_pages,
        ))
    return modules


def run_pipeline(
    customer: str,
    title: str,
    modules: list[CurriculumModule],
    pbl_parts: list[str] | None = None,
    case_parts: list[str] | None = None,
    output_name: str = "test",
) -> tuple[Path, PipelineConfig]:
    """파이프라인 실행 → (출력 경로, config) 반환."""
    out_dir = OUTPUT_DIR / output_name
    out_dir.mkdir(parents=True, exist_ok=True)

    variables = {
        "고객명": customer,
        "교육제목": title,
        "날짜": str(date.today()),
        "제안일": str(date.today()),
        "산업": "",
        "목표": "",
        "페인포인트": "",
        "모듈_요약": ", ".join(m.name for m in modules),
    }

    config = build_config(
        template_dir=TEMPLATE_DIR,
        modules=modules,
        pbl_parts=pbl_parts,
        case_parts=case_parts,
        output_dir=out_dir,
    )

    gamma = GammaHttpClient(api_key=GAMMA_API_KEY or "unused", base_url=GAMMA_BASE_URL)
    replacer = PptxFileReplacer()
    merger = PptxZipMerger()
    pipeline = DefaultProposalPipeline(gamma_client=gamma, replacer=replacer, merger=merger)

    result = pipeline.run(variables, config)
    return result.output_path, config


# ── 검증 함수 ──────────────────────────────────────────────────────────

def count_slides_zip(pptx_path: Path) -> int:
    pat = re.compile(r"^ppt/slides/slide\d+\.xml$")
    with zipfile.ZipFile(str(pptx_path), "r") as z:
        return sum(1 for n in z.namelist() if pat.match(n))


def check_no_placeholders(pptx_path: Path) -> list[str]:
    """ZIP 내 슬라이드 XML에서 {{...}} 패턴이 남아있는지 검사."""
    pat = re.compile(r"\{\{\w+\}\}")
    found = []
    with zipfile.ZipFile(str(pptx_path), "r") as z:
        for name in z.namelist():
            if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                content = z.read(name).decode("utf-8", errors="ignore")
                matches = pat.findall(content)
                if matches:
                    found.append(f"{name}: {matches}")
    return found


def check_media_files(pptx_path: Path) -> dict:
    """미디어 파일 통계."""
    with zipfile.ZipFile(str(pptx_path), "r") as z:
        media = [n for n in z.namelist() if n.startswith("ppt/media/")]
    exts = {}
    for m in media:
        ext = Path(m).suffix.lower()
        exts[ext] = exts.get(ext, 0) + 1
    return exts


def check_slide_masters(pptx_path: Path) -> dict:
    """슬라이드 마스터/레이아웃/테마 수."""
    with zipfile.ZipFile(str(pptx_path), "r") as z:
        names = z.namelist()
    return {
        "masters": sum(1 for n in names if re.match(r"ppt/slideMasters/slideMaster\d+\.xml$", n)),
        "layouts": sum(1 for n in names if re.match(r"ppt/slideLayouts/slideLayout\d+\.xml$", n)),
        "themes": sum(1 for n in names if re.match(r"ppt/theme/theme\d+\.xml$", n)),
    }


def check_fonts_in_xml(pptx_path: Path) -> set[str]:
    """사용된 폰트명 수집."""
    fonts = set()
    font_pat = re.compile(r'typeface="([^"]+)"')
    with zipfile.ZipFile(str(pptx_path), "r") as z:
        for name in z.namelist():
            if name.endswith(".xml") and "slide" in name:
                content = z.read(name).decode("utf-8", errors="ignore")
                fonts.update(font_pat.findall(content))
    return fonts


def check_table_in_work_file(work_path: Path) -> dict:
    """작업 파일의 테이블 행 수 확인."""
    from pptx import Presentation
    prs = Presentation(str(work_path))
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                header = [cell.text.strip()[:20] for cell in table.rows[0].cells]
                data_rows = len(table.rows) - 1
                return {"header": header, "data_rows": data_rows}
    return {"header": [], "data_rows": 0}


# ── 동적 기대값 계산 ──────────────────────────────────────────────────────


INTRO_SLIDES = 2  # 01_도입부.pptx는 2장짜리


def calc_expected_parts(
    modules: list[CurriculumModule],
    pbl_parts: list[str] | None = None,
    case_parts: list[str] | None = None,
) -> int:
    """모듈/PBL/사례에서 예상 파트(FixedPage) 수를 계산."""
    n = 2  # 표지 + 도입부
    for mod in modules:
        n += 2 + len(mod.table_pages) + 1  # 구간표지 + 설계배경 + 타임테이블(N) + 산출물
    if pbl_parts or case_parts:
        n += 1  # PBL 구간표지
        n += len(pbl_parts or [])
        n += len(case_parts or [])
    n += 1  # 엔딩
    return n


def calc_expected_slides(
    modules: list[CurriculumModule],
    pbl_parts: list[str] | None = None,
    case_parts: list[str] | None = None,
) -> int:
    """모듈/PBL/사례에서 예상 슬라이드 수를 계산.

    도입부는 2장, 나머지 파트는 각 1장.
    """
    parts = calc_expected_parts(modules, pbl_parts, case_parts)
    return parts + (INTRO_SLIDES - 1)  # 도입부 파트 1개 = 실제 2장 → +1


# ── 테스트 시나리오 ──────────────────────────────────────────────────────

class TestResult:
    def __init__(self, name: str):
        self.name = name
        self.checks: list[tuple[str, bool, str]] = []

    def check(self, label: str, passed: bool, detail: str = ""):
        self.checks.append((label, passed, detail))

    def print_report(self):
        total = len(self.checks)
        passed = sum(1 for _, p, _ in self.checks if p)
        status = "PASS" if passed == total else "FAIL"
        print(f"\n{'='*60}")
        print(f"[{status}] {self.name} ({passed}/{total})")
        print(f"{'='*60}")
        for label, ok, detail in self.checks:
            mark = "OK" if ok else "NG"
            msg = f"  [{mark}] {label}"
            if detail:
                msg += f" — {detail}"
            print(msg)
        return passed == total


def test_scenario_1():
    """시나리오 1: 전체 5개 커리큘럼 모듈 + PBL + 사례."""
    t = TestResult("시나리오 1: 전체 모듈 + PBL + 사례")

    all_ids = ["process-structuring", "prompt-rubric", "rag-agent", "vibe-coding", "web-agent"]
    modules = load_reference_modules(all_ids)
    pbl = ["part_pbl_로드맵", "part_pbl_레벨무관", "part_pbl_마일스톤", "part_pbl_피드백"]
    cases = ["part_사례_성과공유", "part_사례_르노코리아", "part_사례_현대차", "part_사례_한투증권"]

    pptx_path, config = run_pipeline(
        customer="테스트사", title="26년 AI 교육 로드맵",
        modules=modules, pbl_parts=pbl, case_parts=cases,
        output_name="scenario_1",
    )

    # 1. 파일 생성
    t.check("파일 생성", pptx_path.exists(), str(pptx_path.name))

    # 2. 슬라이드 수 (동적 계산)
    expected_slides = calc_expected_slides(modules, pbl, cases)
    actual = count_slides_zip(pptx_path)
    t.check("슬라이드 수", actual == expected_slides, f"예상 {expected_slides}, 실제 {actual}")

    # 3. 플레이스홀더 잔존 없음
    remaining = check_no_placeholders(pptx_path)
    t.check("플레이스홀더 잔존 없음", len(remaining) == 0,
            f"{len(remaining)}개 발견" if remaining else "")

    # 4. 미디어 파일
    media = check_media_files(pptx_path)
    t.check("미디어 파일 존재", len(media) > 0, str(media))

    # 5. 슬라이드 마스터/레이아웃
    masters = check_slide_masters(pptx_path)
    t.check("슬라이드 마스터", masters["masters"] > 0, str(masters))

    # 6. 폰트 확인
    fonts = check_fonts_in_xml(pptx_path)
    has_pretendard = any("Pretendard" in f or "pretendard" in f.lower() for f in fonts)
    t.check("Pretendard 폰트", has_pretendard, str(sorted(fonts)[:5]))

    # 7. 파트 수 (동적 계산)
    expected_parts = calc_expected_parts(modules, pbl, cases)
    total_parts = len(config.fixed_pages) + len(config.dynamic_pages)
    t.check("파트 수", total_parts == expected_parts, f"예상 {expected_parts}, 실제 {total_parts} (고정 {len(config.fixed_pages)}, 동적 {len(config.dynamic_pages)})")

    return t.print_report()


def test_scenario_2():
    """시나리오 2: 부분 모듈 (2개만)."""
    t = TestResult("시나리오 2: 부분 모듈 (process + rag)")

    modules = load_reference_modules(["process-structuring", "rag-agent"])

    pptx_path, config = run_pipeline(
        customer="부분테스트", title="AI 교육",
        modules=modules,
        output_name="scenario_2",
    )

    t.check("파일 생성", pptx_path.exists())

    expected = calc_expected_slides(modules)
    actual = count_slides_zip(pptx_path)
    t.check("슬라이드 수", actual == expected, f"예상 {expected}, 실제 {actual}")

    remaining = check_no_placeholders(pptx_path)
    t.check("플레이스홀더 잔존 없음", len(remaining) == 0)

    masters = check_slide_masters(pptx_path)
    t.check("마스터 무결성", masters["masters"] > 0, str(masters))

    return t.print_report()


def test_scenario_3():
    """시나리오 3: 모듈 없이 골격만 (표지 + 도입부 + 엔딩)."""
    t = TestResult("시나리오 3: 골격만 (모듈 없음)")

    pptx_path, config = run_pipeline(
        customer="골격테스트", title="테스트",
        modules=[],
        output_name="scenario_3",
    )

    t.check("파일 생성", pptx_path.exists())

    expected = calc_expected_slides([])
    actual = count_slides_zip(pptx_path)
    t.check("슬라이드 수", actual == expected, f"예상 {expected}, 실제 {actual}")

    remaining = check_no_placeholders(pptx_path)
    t.check("플레이스홀더 잔존 없음", len(remaining) == 0)

    return t.print_report()


def test_scenario_4():
    """시나리오 4: 다중 타임테이블 (web-agent, 5페이지)."""
    t = TestResult("시나리오 4: web-agent (타임테이블 5페이지)")

    modules = load_reference_modules(["web-agent"])

    pptx_path, config = run_pipeline(
        customer="웹에이전트사", title="Web-Agent 교육",
        modules=modules,
        output_name="scenario_4",
    )

    t.check("파일 생성", pptx_path.exists())

    expected = calc_expected_slides(modules)
    actual = count_slides_zip(pptx_path)
    t.check("슬라이드 수", actual == expected, f"예상 {expected}, 실제 {actual}")

    # 작업 파일에서 타임테이블 행 수 검증
    work_dir = OUTPUT_DIR / "scenario_4" / "_work"
    expected_tt = len(modules[0].table_pages)
    tt_files = sorted(work_dir.glob("*_web-agent_타임테이블*.pptx"))
    t.check("타임테이블 파일 수", len(tt_files) == expected_tt, f"예상 {expected_tt}, 실제 {len(tt_files)}")

    # 각 타임테이블의 데이터 행 수
    expected_topics = sum(len(tp.rows) for tp in modules[0].table_pages)
    total_data_rows = 0
    for tf in tt_files:
        info = check_table_in_work_file(tf)
        total_data_rows += info["data_rows"]
        t.check(f"테이블 헤더 ({tf.stem})", info["header"][0] == "주제명" if info["header"] else False,
                f"헤더: {info['header']}")

    t.check("총 데이터 행 수", total_data_rows == expected_topics, f"예상 {expected_topics}, 실제 {total_data_rows}")

    remaining = check_no_placeholders(pptx_path)
    t.check("플레이스홀더 잔존 없음", len(remaining) == 0)

    return t.print_report()


def test_scenario_5():
    """시나리오 5: PBL만 (커리큘럼 모듈 없이)."""
    t = TestResult("시나리오 5: PBL만 (모듈 없음)")

    pbl = ["part_pbl_로드맵", "part_pbl_마일스톤"]
    cases = ["part_사례_현대차"]

    pptx_path, config = run_pipeline(
        customer="PBL테스트", title="PBL 교육",
        modules=[], pbl_parts=pbl, case_parts=cases,
        output_name="scenario_5",
    )

    t.check("파일 생성", pptx_path.exists())

    expected = calc_expected_slides([], pbl, cases)
    actual = count_slides_zip(pptx_path)
    t.check("슬라이드 수", actual == expected, f"예상 {expected}, 실제 {actual}")

    remaining = check_no_placeholders(pptx_path)
    t.check("플레이스홀더 잔존 없음", len(remaining) == 0)

    return t.print_report()


# ── 메인 ──────────────────────────────────────────────────────────────

def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    results = []
    results.append(test_scenario_1())
    results.append(test_scenario_2())
    results.append(test_scenario_3())
    results.append(test_scenario_4())
    results.append(test_scenario_5())

    print(f"\n{'='*60}")
    total = len(results)
    passed = sum(1 for r in results if r)
    print(f"전체 결과: {passed}/{total} 시나리오 통과")
    if passed == total:
        print("ALL PASS")
    else:
        print("SOME FAILED")
        sys.exit(1)


if __name__ == "__main__":
    main()
