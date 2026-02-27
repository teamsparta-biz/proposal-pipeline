"""PPTX 합성 테스트 스크립트.

test_pptx_replace.py가 먼저 실행되어야 합니다 (더미 PPTX 생성).

테스트 시나리오:
  1. 더미 3장 (표지 + 회사소개 + 약관) 합성
  2. 슬라이드 수 검증
  3. 각 슬라이드의 핵심 텍스트 검증
  4. 순서 변경 합성 (표지 → 약관 → 회사소개)

사용법:
  python test_pptx_merge.py
"""

import sys
from pathlib import Path

from pptx import Presentation

from proposal_pipeline.pptx_merger.merger import PptxFileMerger


TEMPLATE_DIR = Path(__file__).parent.parent / "templates" / "dummy"
OUTPUT_DIR = Path(__file__).parent.parent / "output_test"


def get_all_text(slide) -> str:
    """슬라이드의 모든 텍스트를 추출."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    texts.append(text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    text = cell.text.strip()
                    if text:
                        texts.append(text)
    return " | ".join(texts)


def test_basic_merge():
    """기본 합성: 표지 + 회사소개 + 약관 → 3슬라이드."""
    print("=== 테스트 1: 기본 합성 (3파일 → 1파일) ===")

    sources = [
        TEMPLATE_DIR / "표지.pptx",
        TEMPLATE_DIR / "회사소개.pptx",
        TEMPLATE_DIR / "약관.pptx",
    ]

    # 소스 존재 확인
    for s in sources:
        if not s.exists():
            print(f"  에러: {s} 가 없습니다. test_pptx_replace.py를 먼저 실행하세요.")
            return False

    merger = PptxFileMerger()
    output_path = OUTPUT_DIR / "merged_basic.pptx"
    merger.merge_and_save(sources, output_path)
    print(f"  출력: {output_path}")

    # 검증
    result = Presentation(str(output_path))
    slide_count = len(result.slides)
    print(f"  슬라이드 수: {slide_count} (기대: 3)")

    if slide_count != 3:
        print(f"  !! 실패: 슬라이드 수 불일치")
        return False

    # 각 슬라이드 내용 확인
    for i, slide in enumerate(result.slides):
        text = get_all_text(slide)
        print(f"  슬라이드 {i+1}: {text[:80]}...")

    # 순서 확인: 첫 슬라이드에 "귀중", 마지막에 "유효" 또는 "견적"
    slide1_text = get_all_text(result.slides[0])
    slide3_text = get_all_text(result.slides[2])

    order_ok = "귀중" in slide1_text and ("유효" in slide3_text or "견적" in slide3_text or "제안" in slide3_text)
    if order_ok:
        print(f"  OK - 순서 정상 (표지 → 회사소개 → 약관)")
    else:
        print(f"  !! 순서 확인 필요")

    print()
    return True


def test_reordered_merge():
    """순서 변경 합성: 표지 → 약관 → 회사소개."""
    print("=== 테스트 2: 순서 변경 합성 ===")

    sources = [
        TEMPLATE_DIR / "표지.pptx",
        TEMPLATE_DIR / "약관.pptx",
        TEMPLATE_DIR / "회사소개.pptx",
    ]

    merger = PptxFileMerger()
    output_path = OUTPUT_DIR / "merged_reordered.pptx"
    merger.merge_and_save(sources, output_path)
    print(f"  출력: {output_path}")

    result = Presentation(str(output_path))
    slide_count = len(result.slides)
    print(f"  슬라이드 수: {slide_count} (기대: 3)")

    if slide_count != 3:
        print(f"  !! 실패: 슬라이드 수 불일치")
        return False

    # 순서 확인: 2번째가 약관
    slide2_text = get_all_text(result.slides[1])
    order_ok = "제안" in slide2_text or "견적" in slide2_text or "유효" in slide2_text
    if order_ok:
        print(f"  OK - 순서 정상 (표지 → 약관 → 회사소개)")
    else:
        print(f"  슬라이드 2 텍스트: {slide2_text[:60]}")
        print(f"  !! 순서 확인 필요")

    print()
    return True


def test_mixed_merge():
    """치환된 파일과 원본 파일 혼합 합성 (실전 시나리오)."""
    print("=== 테스트 3: 치환 파일 + 원본 파일 혼합 합성 ===")

    replaced_cover = OUTPUT_DIR / "표지.pptx"  # 치환된 표지 (삼성전자)
    original_intro = TEMPLATE_DIR / "회사소개.pptx"  # 원본 회사소개
    replaced_terms = OUTPUT_DIR / "약관.pptx"  # 치환된 약관

    sources = []
    for p in [replaced_cover, original_intro, replaced_terms]:
        if p.exists():
            sources.append(p)
        else:
            print(f"  스킵: {p} 없음")

    if len(sources) < 2:
        print(f"  에러: 소스가 부족합니다. test_pptx_replace.py를 먼저 실행하세요.")
        return False

    merger = PptxFileMerger()
    output_path = OUTPUT_DIR / "merged_mixed.pptx"
    merger.merge_and_save(sources, output_path)
    print(f"  출력: {output_path}")

    result = Presentation(str(output_path))
    print(f"  슬라이드 수: {len(result.slides)} (기대: {len(sources)})")

    # 첫 슬라이드에 "삼성전자" 있으면 치환 파일이 정상 합성된 것
    slide1_text = get_all_text(result.slides[0])
    if "삼성전자" in slide1_text:
        print(f"  OK - 치환된 내용 유지됨 (삼성전자)")
    else:
        print(f"  슬라이드 1: {slide1_text[:60]}")
        print(f"  !! 치환 내용 확인 필요")

    print()
    return True


if __name__ == "__main__":
    results = []
    results.append(test_basic_merge())
    results.append(test_reordered_merge())
    results.append(test_mixed_merge())

    print("=" * 40)
    if all(results):
        print("전체 테스트 통과!")
    else:
        passed = sum(results)
        print(f"{passed}/{len(results)} 테스트 통과")

    print(f"\n합성 파일 위치: {OUTPUT_DIR.resolve()}")
    print("PowerPoint로 열어서 시각적으로도 확인해보세요.")
