"""PPTX 치환 엔진 테스트 스크립트.

1단계: 더미 PPTX 3장 생성 (표지, 회사소개, 약관)
2단계: 플레이스홀더 치환 실행
3단계: 결과 검증 (텍스트 비교)

사용법:
  pip install python-pptx
  python test_pptx_replace.py
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from proposal_pipeline.pptx_replacer.replacer import PptxFileReplacer


TEMPLATE_DIR = Path(__file__).parent.parent / "templates" / "dummy"
OUTPUT_DIR = Path(__file__).parent.parent / "output_test"


# === 1단계: 더미 PPTX 생성 ===

def create_cover():
    """표지 슬라이드 — 플레이스홀더: 고객명, 날짜."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # 제목
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "{{고객명}} 귀중"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # 부제
    txBox2 = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(9), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "AI 업무 자동화 교육 제안서"
    run2.font.size = Pt(24)
    run2.font.color.rgb = RGBColor(0x44, 0x44, 0x66)

    # 날짜 + 회사명
    txBox3 = slide.shapes.add_textbox(Inches(3), Inches(5.5), Inches(7), Inches(0.6))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "팀스파르타 | {{날짜}}"
    run3.font.size = Pt(16)
    run3.font.color.rgb = RGBColor(0x88, 0x88, 0x99)

    path = TEMPLATE_DIR / "표지.pptx"
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    print(f"  생성: {path}")
    return path


def create_company_intro():
    """회사소개 슬라이드 — 플레이스홀더: 담당자명, 연락처."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "회사 소개"
    run.font.size = Pt(28)
    run.font.bold = True

    # 본문
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(4))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    lines = [
        "팀스파르타는 AI 교육 전문 기업입니다.",
        "",
        "주요 서비스:",
        "  - AI Literacy 교육",
        "  - 맞춤형 업무 자동화 컨설팅",
        "  - PBL 기반 실습 프로그램",
        "",
        "담당자: {{담당자명}}",
        "연락처: {{연락처}}",
    ]
    for i, line in enumerate(lines):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(16)

    path = TEMPLATE_DIR / "회사소개.pptx"
    prs.save(str(path))
    print(f"  생성: {path}")
    return path


def create_terms():
    """약관 슬라이드 — 플레이스홀더: 제안일, 금액, 고객명. 테이블 포함."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "제안 조건"
    run.font.size = Pt(28)
    run.font.bold = True

    # 테이블
    rows, cols = 4, 2
    table_shape = slide.shapes.add_table(rows, cols, Inches(1.5), Inches(2), Inches(10), Inches(3))
    table = table_shape.table

    data = [
        ("항목", "내용"),
        ("제안 대상", "{{고객명}}"),
        ("유효 기간", "{{제안일}}부터 30일"),
        ("견적 금액", "{{금액}}원 (VAT 별도)"),
    ]
    for r, (c1, c2) in enumerate(data):
        table.cell(r, 0).text = c1
        table.cell(r, 1).text = c2

    path = TEMPLATE_DIR / "약관.pptx"
    prs.save(str(path))
    print(f"  생성: {path}")
    return path


# === 2단계 + 3단계: 치환 + 검증 ===

def test_replace():
    variables = {
        "고객명": "삼성전자",
        "날짜": "2026-02-25",
        "담당자명": "김철수",
        "연락처": "010-1234-5678",
        "제안일": "2026-02-25",
        "금액": "30,000,000",
    }

    replacer = PptxFileReplacer()
    files = ["표지.pptx", "회사소개.pptx", "약관.pptx"]

    print()
    print("=== 치환 실행 ===")

    all_passed = True
    for fname in files:
        src = TEMPLATE_DIR / fname
        dst = OUTPUT_DIR / fname
        replacer.replace_and_save(src, variables, dst)
        print(f"  치환 완료: {dst}")

        # 검증: 치환된 파일을 다시 읽어서 플레이스홀더가 남아있는지 확인
        prs = Presentation(str(dst))
        remaining = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs)
                        if "{{" in text:
                            remaining.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text = cell.text
                            if "{{" in text:
                                remaining.append(text)

        if remaining:
            print(f"    !! 미치환 플레이스홀더 발견: {remaining}")
            all_passed = False
        else:
            print(f"    OK - 플레이스홀더 모두 치환됨")

    # 내용 확인: 표지에서 고객명 확인
    print()
    print("=== 내용 확인 ===")
    cover = Presentation(str(OUTPUT_DIR / "표지.pptx"))
    for shape in cover.slides[0].shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs)
                if text.strip():
                    print(f"  {text}")

    terms = Presentation(str(OUTPUT_DIR / "약관.pptx"))
    for shape in terms.slides[0].shapes:
        if shape.has_table:
            print()
            print("  [약관 테이블]")
            for row in shape.table.rows:
                cells = [cell.text for cell in row.cells]
                print(f"  | {cells[0]:12s} | {cells[1]}")

    print()
    if all_passed:
        # 서식 보존 확인: 표지 제목의 볼드/크기
        for shape in cover.slides[0].shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "삼성전자" in run.text:
                            bold_ok = run.font.bold == True
                            size_ok = run.font.size == Pt(36)
                            print(f"  서식 확인: bold={run.font.bold} (기대: True), size={run.font.size} (기대: {Pt(36)})")
                            if not (bold_ok and size_ok):
                                all_passed = False
                                print(f"    !! 서식이 변경됨")
                            else:
                                print(f"    OK - 서식 유지됨")

    print()
    if all_passed:
        print("=== 전체 테스트 통과! ===")
    else:
        print("=== 일부 테스트 실패 ===")

    print(f"\n치환된 파일 위치: {OUTPUT_DIR.resolve()}")
    print("PowerPoint로 열어서 시각적으로도 확인해보세요.")


if __name__ == "__main__":
    print("=== 1단계: 더미 PPTX 생성 ===")
    create_cover()
    create_company_intro()
    create_terms()

    test_replace()
