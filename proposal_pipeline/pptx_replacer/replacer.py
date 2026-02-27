"""PPTX 플레이스홀더 치환 + 테이블 주입 엔진 — ABC 인터페이스 + 구현체."""

from __future__ import annotations

import copy
import re
from abc import ABC, abstractmethod
from io import BytesIO
from pathlib import Path
from typing import Optional

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from ..curriculum.models import TopicRow


class PptxReplacer(ABC):
    """PPTX 플레이스홀더 치환 추상 인터페이스."""

    @abstractmethod
    def replace(self, pptx_path: Path, variables: dict[str, str]) -> Presentation:
        """PPTX 파일의 {{변수}}를 치환하여 Presentation 객체를 반환."""

    @abstractmethod
    def replace_and_save(self, pptx_path: Path, variables: dict[str, str], output_path: Path) -> Path:
        """치환 후 파일로 저장."""

    @abstractmethod
    def inject_table(self, prs: Presentation, slide_index: int, rows: list[TopicRow]) -> None:
        """슬라이드의 테이블에 TopicRow 데이터를 주입한다.

        헤더 행(0행)은 유지하고, 데이터 행(1행~)을 교체한다.
        행 수가 다르면 추가/삭제하여 맞춘다.
        """

    @abstractmethod
    def replace_slide_image(self, prs: Presentation, slide_index: int, image_path: Path) -> None:
        """슬라이드에서 가장 큰 PICTURE shape의 이미지를 교체한다."""


class PptxFileReplacer(PptxReplacer):
    """python-pptx 기반 플레이스홀더 치환 구현체."""

    PLACEHOLDER_PATTERN = re.compile(r"\{\{(\w+)\}\}")

    def replace(self, pptx_path: Path, variables: dict[str, str]) -> Presentation:
        prs = Presentation(str(pptx_path))
        for slide in prs.slides:
            self._replace_slide(slide, variables)
        return prs

    def replace_and_save(self, pptx_path: Path, variables: dict[str, str], output_path: Path) -> Path:
        prs = self.replace(pptx_path, variables)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        return output_path

    def _replace_slide(self, slide, variables: dict[str, str]):
        """슬라이드 내 모든 shape의 텍스트를 치환."""
        for shape in slide.shapes:
            # 텍스트 프레임
            if shape.has_text_frame:
                self._replace_text_frame(shape.text_frame, variables)

            # 테이블
            if shape.has_table:
                self._replace_table(shape.table, variables)

    def _replace_text_frame(self, text_frame, variables: dict[str, str]):
        """텍스트 프레임 내 모든 paragraph의 플레이스홀더를 치환.

        PowerPoint는 텍스트를 여러 run으로 쪼갤 수 있다.
        예: "{{고객명}}"이 "{{고", "객명", "}}" 세 개의 run으로 나뉠 수 있음.
        이를 해결하기 위해 paragraph 단위로 전체 텍스트를 조합한 후 치환하고,
        첫 번째 run에 결과를 넣고 나머지 run을 비운다.
        """
        for paragraph in text_frame.paragraphs:
            full_text = "".join(run.text for run in paragraph.runs)

            if not self.PLACEHOLDER_PATTERN.search(full_text):
                continue

            # 치환 수행
            new_text = full_text
            for key, value in variables.items():
                new_text = new_text.replace(f"{{{{{key}}}}}", str(value))

            # run이 없으면 스킵
            if not paragraph.runs:
                continue

            # 첫 번째 run에 전체 텍스트 할당 (서식 유지)
            paragraph.runs[0].text = new_text
            # 나머지 run 비우기
            for run in paragraph.runs[1:]:
                run.text = ""

    def _replace_table(self, table, variables: dict[str, str]):
        """테이블 셀 내 플레이스홀더를 치환."""
        for row in table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    self._replace_text_frame(cell.text_frame, variables)

    # ── 테이블 주입 ──

    def inject_table(self, prs: Presentation, slide_index: int, rows: list[TopicRow]) -> None:
        """슬라이드의 테이블에 TopicRow 데이터를 주입한다."""
        slide = prs.slides[slide_index]
        table = self._find_table(slide)
        if table is None:
            raise ValueError(f"슬라이드 {slide_index}에 테이블이 없습니다")

        tbl_xml = table._tbl  # lxml Element

        # 기존 데이터 행 수집 (0행 = 헤더, 1행~ = 데이터)
        existing_rows = list(tbl_xml.iterchildren(qn("a:tr")))
        header_row = existing_rows[0]
        data_rows = existing_rows[1:]

        # 서식 템플릿: 첫 번째 데이터 행을 깊은 복사하여 보존
        if not data_rows:
            raise ValueError("테이블에 데이터 행이 없어 서식 템플릿을 생성할 수 없습니다")
        template_row = copy.deepcopy(data_rows[0])

        # 기존 데이터 행 모두 제거
        for row_el in data_rows:
            tbl_xml.remove(row_el)

        # 새 데이터 행 추가
        col_fields = ["subject", "hours", "content", "exercise"]
        for topic in rows:
            new_row = copy.deepcopy(template_row)
            cells = list(new_row.iterchildren(qn("a:tc")))
            for ci, field_name in enumerate(col_fields):
                if ci < len(cells):
                    self._set_cell_text(cells[ci], getattr(topic, field_name))
            tbl_xml.append(new_row)

        # 행 높이 자동 조정 (테이블 오버플로 방지)
        data_row_height = self._auto_fit_rows(prs, slide, len(rows))

        # 셀 폰트 크기 자동 축소 (텍스트가 행 높이를 초과하면)
        if data_row_height:
            table_shape = None
            for shape in slide.shapes:
                if shape.has_table:
                    table_shape = shape
                    break
            if table_shape:
                self._auto_fit_cell_font(table_shape, data_row_height)

    @staticmethod
    def _auto_fit_rows(prs: Presentation, slide, num_data_rows: int) -> int | None:
        """데이터 행 높이를 슬라이드에 맞게 자동 조정한다.

        슬라이드 높이에서 테이블 상단 위치와 하단 마진을 빼고,
        헤더 행을 제외한 나머지 공간을 데이터 행 수로 균등 분배한다.

        Returns:
            설정된 데이터 행 높이 (EMU), 또는 조정 불필요 시 None.
        """
        if num_data_rows <= 0:
            return None

        # 슬라이드에서 테이블 shape 찾기 (top 위치 필요)
        table_shape = None
        for shape in slide.shapes:
            if shape.has_table:
                table_shape = shape
                break
        if table_shape is None:
            return None

        slide_height = prs.slide_height
        table_top = table_shape.top
        bottom_margin = 274320  # 0.3in in EMU
        min_row_height = 457200  # 0.5in in EMU

        available_height = slide_height - table_top - bottom_margin

        # 헤더 행 높이 구하기
        tbl_xml = table_shape.table._tbl
        tr_elements = list(tbl_xml.iterchildren(qn("a:tr")))
        if len(tr_elements) < 2:
            return None

        header_tr = tr_elements[0]
        header_h = int(header_tr.get("h", "0"))

        # 데이터 행에 할당 가능한 높이
        data_area = available_height - header_h
        data_row_height = max(data_area // num_data_rows, min_row_height)

        # 데이터 행(헤더 이후)에 높이 설정
        for tr in tr_elements[1:]:
            tr.set("h", str(data_row_height))

        return data_row_height

    # 폰트 축소 대상 컬럼 인덱스 (content=2, exercise=3)
    _DENSE_COLS = {2, 3}

    @staticmethod
    def _auto_fit_cell_font(table_shape, data_row_height_emu: int) -> None:
        """content·exercise 셀의 폰트를 행 높이에 맞게 자동 축소한다.

        텍스트가 많은 컬럼(학습내용, 실습 예시)만 대상으로 하여,
        주제명·시간 컬럼의 가독성은 유지한다.
        """
        tbl = table_shape.table._tbl

        # 컬럼 너비 (EMU)
        grid = tbl.find(qn("a:tblGrid"))
        col_widths = [int(gc.get("w", "0")) for gc in grid.iterchildren(qn("a:gridCol"))]

        tr_elements = list(tbl.iterchildren(qn("a:tr")))
        if len(tr_elements) < 2:
            return

        # 기본 폰트 크기 추출 (첫 데이터 행 → 첫 셀)
        base_sz = 1100  # default 11pt (hundredths of a point)
        first_data = tr_elements[1]
        for tc in first_data.iterchildren(qn("a:tc")):
            txBody = tc.find(qn("a:txBody"))
            if txBody is None:
                continue
            for p in txBody.iterchildren(qn("a:p")):
                for r in p.iterchildren(qn("a:r")):
                    rPr = r.find(qn("a:rPr"))
                    if rPr is not None and rPr.get("sz"):
                        base_sz = int(rPr.get("sz"))
                        break
                break
            break

        font_pt = base_sz / 100  # pt 단위

        # 대상 컬럼에서만 행별 최대 줄 수 계산
        worst_lines = 0
        for tr in tr_elements[1:]:
            cells = list(tr.iterchildren(qn("a:tc")))
            for ci, tc in enumerate(cells):
                if ci not in PptxFileReplacer._DENSE_COLS:
                    continue
                if ci >= len(col_widths):
                    break
                text = PptxFileReplacer._extract_cell_text(tc)
                if not text:
                    continue
                lines = PptxFileReplacer._estimate_text_lines(
                    text, col_widths[ci], font_pt
                )
                worst_lines = max(worst_lines, lines)

        if worst_lines <= 0:
            return

        # 행 높이 → 가용 높이 (pt)  [1 inch = 914400 EMU = 72 pt]
        row_h_pt = data_row_height_emu * 72 / 914400
        cell_padding_pt = 12  # 상하 내부 패딩 합계
        usable_pt = row_h_pt - cell_padding_pt

        # 현재 폰트에서 필요한 높이
        line_height_pt = font_pt * 1.3  # 줄간격 포함
        needed_pt = worst_lines * line_height_pt

        if needed_pt <= usable_pt:
            return  # 축소 불필요

        # 축소 비율 계산
        scale = usable_pt / needed_pt
        new_sz = max(int(base_sz * scale), 850)  # 최소 8.5pt

        # 대상 컬럼에만 적용
        for tr in tr_elements[1:]:
            cells = list(tr.iterchildren(qn("a:tc")))
            for ci, tc in enumerate(cells):
                if ci in PptxFileReplacer._DENSE_COLS:
                    PptxFileReplacer._set_cell_font_size(tc, new_sz)

    @staticmethod
    def _extract_cell_text(tc_element) -> str:
        """테이블 셀에서 전체 텍스트를 추출한다."""
        txBody = tc_element.find(qn("a:txBody"))
        if txBody is None:
            return ""
        parts = []
        for p in txBody.iterchildren(qn("a:p")):
            for r in p.iterchildren(qn("a:r")):
                t = r.find(qn("a:t"))
                if t is not None and t.text:
                    parts.append(t.text)
        return "".join(parts)

    @staticmethod
    def _estimate_text_lines(text: str, col_width_emu: int, font_pt: float) -> int:
        """텍스트가 셀에서 차지하는 줄 수를 추정한다.

        한글/영문 혼합 기준으로 문자 너비를 추정하고,
        \\n 줄바꿈과 자동 줄바꿈을 모두 고려한다.
        """
        import math

        # 셀 너비 → 인치 → 가용 너비 (좌우 패딩 차감)
        col_inches = col_width_emu / 914400
        usable_inches = max(col_inches - 0.2, 0.5)

        # 한글+영문 혼합 기준 문자 너비: font_pt * 0.012 인치/글자 (경험적 값)
        char_width_inches = font_pt * 0.012
        chars_per_line = max(1, int(usable_inches / char_width_inches))

        # \n 으로 분할하여 각 줄의 줄바꿈 계산
        total_lines = 0
        for line in text.split("\n"):
            line = line.strip()
            if not line:
                total_lines += 0.5  # 빈 줄은 반줄
            else:
                total_lines += math.ceil(len(line) / chars_per_line)

        return int(math.ceil(total_lines))

    @staticmethod
    def _set_cell_font_size(tc_element, size_hundredths: int) -> None:
        """테이블 셀의 모든 run에 폰트 크기를 설정한다."""
        txBody = tc_element.find(qn("a:txBody"))
        if txBody is None:
            return
        for p in txBody.iterchildren(qn("a:p")):
            for r in p.iterchildren(qn("a:r")):
                rPr = r.find(qn("a:rPr"))
                if rPr is not None:
                    rPr.set("sz", str(size_hundredths))

    def replace_slide_image(self, prs: Presentation, slide_index: int, image_path: Path) -> None:
        """슬라이드에서 가장 큰 PICTURE shape의 이미지 blob을 교체한다."""
        slide = prs.slides[slide_index]
        image_path = Path(image_path)

        # 가장 큰 PICTURE shape 찾기 (면적 기준)
        best_shape = None
        best_area = 0
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                area = shape.width * shape.height
                if area > best_area:
                    best_area = area
                    best_shape = shape

        if best_shape is None:
            raise ValueError(f"슬라이드 {slide_index}에 PICTURE shape가 없습니다")

        # 이미지 파트의 blob 교체
        blip_elements = best_shape._element.findall(".//" + qn("a:blip"))
        if not blip_elements:
            raise ValueError(f"shape '{best_shape.name}'에 blip 요소가 없습니다")

        r_id = blip_elements[0].get(qn("r:embed"))
        image_part = slide.part.rels[r_id].target_part
        image_part.blob = image_path.read_bytes()

    @staticmethod
    def _find_table(slide):
        """슬라이드에서 첫 번째 테이블 shape를 반환한다."""
        for shape in slide.shapes:
            if shape.has_table:
                return shape.table
        return None

    @staticmethod
    def _sanitize_text(text: str) -> str:
        """XML 호환 텍스트로 정제한다. 제어 문자를 줄바꿈으로 변환."""
        # vertical tab (\x0b), form feed (\x0c) 등 → 줄바꿈
        text = text.replace("\x0b", "\n").replace("\x0c", "\n")
        # 나머지 XML 비호환 제어 문자 제거 (탭/줄바꿈 제외)
        return re.sub(r"[\x00-\x08\x0e-\x1f]", "", text)

    @staticmethod
    def _set_cell_text(tc_element, text: str):
        """테이블 셀(a:tc)의 텍스트를 교체한다. 첫 번째 run의 서식을 유지."""
        # 제어 문자 정제
        text = PptxFileReplacer._sanitize_text(text)

        paragraphs = list(tc_element.iterchildren(qn("a:txBody")))
        if not paragraphs:
            return
        tx_body = paragraphs[0]

        p_elements = list(tx_body.iterchildren(qn("a:p")))
        if not p_elements:
            return

        # 첫 번째 paragraph의 첫 run에 전체 텍스트 삽입
        first_p = p_elements[0]
        runs = list(first_p.iterchildren(qn("a:r")))
        if runs:
            t_el = runs[0].find(qn("a:t"))
            if t_el is not None:
                t_el.text = text
            for r in runs[1:]:
                first_p.remove(r)
        # 나머지 paragraph 제거
        for p in p_elements[1:]:
            tx_body.remove(p)
